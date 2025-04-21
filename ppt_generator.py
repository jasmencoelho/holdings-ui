from pptx import Presentation
from pptx.util import Inches
import os

def create_slide_with_images(prs, title_text, image_paths):
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Title at the top
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title_text

    # Fixed positions: space them evenly across the slide (side-by-side)
    img_width = Inches(2.5)
    img_height = Inches(2.5)
    img_top = Inches(1.0)
    img_spacing = Inches(0.5)
    left_start = Inches(1)

    for i, path in enumerate(image_paths):
        left = left_start + i * (img_width + img_spacing)
        slide.shapes.add_picture(path, left, img_top, width=img_width, height=img_height)

def build_ppt_for_regions(region_images_dict, out_path):
    import os
    from pptx import Presentation
    from pptx.util import Inches

    os.makedirs(os.path.dirname(out_path), exist_ok=True)
    prs = Presentation()

    for region_code, image_paths in region_images_dict.items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
        title_frame = title_box.text_frame
        title_frame.text = region_code

        # Place images at fixed size and position
        img_width = Inches(2.5)
        img_height = Inches(2.5)
        img_top = Inches(1.0)
        img_spacing = Inches(0.5)
        left_start = Inches(1)

        for i, path in enumerate(image_paths):
            left = left_start + i * (img_width + img_spacing)
            slide.shapes.add_picture(path, left, img_top, width=img_width, height=img_height)

    prs.save(out_path)
